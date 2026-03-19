import html
import os
import re
import shutil
import subprocess
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import List

import docx2txt
from llama_index.core import Document, SimpleDirectoryReader
from openpyxl import load_workbook
from pptx import Presentation

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".pptx",
    ".xlsx",
    ".txt",
    ".md",
    ".html",
    ".jpg",
    ".jpeg",
    ".png",
    ".xdmp",
}


class DocumentLoader:
    """Unified loader for supported office, text, and fallback file formats."""

    def load(self, file_path: str) -> List[Document]:
        ext = Path(file_path).suffix.lower()

        if ext == ".pdf":
            return self._load_pdf(file_path)
        if ext == ".docx":
            return self._load_docx(file_path)
        if ext == ".doc":
            return self._load_doc(file_path)
        if ext == ".pptx":
            return self._load_pptx(file_path)
        if ext == ".xlsx":
            return self._load_xlsx(file_path)
        if ext in {".txt", ".md", ".xdmp"}:
            return self._load_text(file_path)
        if ext == ".html":
            return self._load_html(file_path)
        if ext in {".jpg", ".jpeg", ".png"}:
            return self._load_image(file_path)

        raise ValueError(f"Unsupported file type: {ext or '<none>'}")

    def load_many(self, file_paths: List[str]) -> List[Document]:
        documents: List[Document] = []
        for file_path in file_paths:
            documents.extend(self.load(file_path))
        return documents

    def _base_metadata(self, file_path: str) -> dict:
        name = os.path.basename(file_path)
        return {
            "file_name": name,
            "source": name,
            "path": os.path.realpath(file_path),
        }

    def _make_document(self, text: str, file_path: str) -> List[Document]:
        content = (text or "").strip()
        if not content:
            raise ValueError(f"No extractable text found in {os.path.basename(file_path)}")
        return [Document(text=content, metadata=self._base_metadata(file_path))]

    def _load_with_simple_reader(self, file_path: str) -> List[Document]:
        reader = SimpleDirectoryReader(input_files=[file_path])
        documents = reader.load_data()
        normalized: List[Document] = []
        base_metadata = self._base_metadata(file_path)
        for doc in documents:
            metadata = getattr(doc, "metadata", {}) or {}
            metadata = {**base_metadata, **metadata}
            doc.metadata = metadata
            normalized.append(doc)
        return normalized

    def _load_pdf(self, file_path: str) -> List[Document]:
        return self._load_with_simple_reader(file_path)

    def _load_image(self, file_path: str) -> List[Document]:
        return self._load_with_simple_reader(file_path)

    def _read_text(self, file_path: str) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def _load_text(self, file_path: str) -> List[Document]:
        return self._make_document(self._read_text(file_path), file_path)

    def _load_html(self, file_path: str) -> List[Document]:
        raw = self._read_text(file_path)
        text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", raw)
        text = re.sub(r"(?s)<[^>]+>", " ", text)
        text = html.unescape(text)
        text = re.sub(r"\s+", " ", text).strip()
        return self._make_document(text, file_path)

    def _load_docx(self, file_path: str) -> List[Document]:
        text = docx2txt.process(file_path) or ""
        return self._make_document(text, file_path)

    def _load_doc(self, file_path: str) -> List[Document]:
        converted_text = self._convert_doc_to_docx_text(file_path)
        if converted_text:
            return self._make_document(converted_text, file_path)

        extracted = self._extract_doc_with_legacy_tools(file_path)
        if extracted:
            return self._make_document(extracted, file_path)

        raise ValueError(
            "Legacy .doc parsing requires LibreOffice (`soffice`) or a compatible text extractor."
        )

    def _convert_doc_to_docx_text(self, file_path: str) -> str | None:
        soffice = shutil.which("soffice")
        if not soffice:
            return None

        with TemporaryDirectory() as output_dir:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "docx",
                    file_path,
                    "--outdir",
                    output_dir,
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
            converted = os.path.join(
                output_dir,
                f"{Path(file_path).stem}.docx",
            )
            if not os.path.exists(converted):
                return None
            return docx2txt.process(converted) or None

    def _extract_doc_with_legacy_tools(self, file_path: str) -> str | None:
        for cmd in ("antiword", "catdoc"):
            tool = shutil.which(cmd)
            if not tool:
                continue
            proc = subprocess.run(
                [tool, file_path],
                check=False,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="ignore",
            )
            text = (proc.stdout or "").strip()
            if proc.returncode == 0 and text:
                return text
        return None

    def _load_pptx(self, file_path: str) -> List[Document]:
        prs = Presentation(file_path)
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    texts.append(text)
        return self._make_document("\n".join(texts), file_path)

    def _load_xlsx(self, file_path: str) -> List[Document]:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        lines: List[str] = []
        for sheet in wb.worksheets:
            lines.append(f"[Sheet] {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if values:
                    lines.append("\t".join(values))
        return self._make_document("\n".join(lines), file_path)
